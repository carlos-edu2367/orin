"""Native text extraction — no model is involved here.

Most real PDFs carry a text layer, and every Office document does. Reading them
directly is both free and better than any transcription, so the visual path is
only ever reached for what genuinely has no text.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

MAX_EXTRACTED_CHARS = 40_000
MAX_PDF_PAGES = 200


@dataclass(frozen=True, slots=True)
class ExtractedText:
    text: str
    truncated: bool = False
    # 1-based page numbers of a PDF whose page carried no text layer.
    pages_without_text: tuple[int, ...] = ()


def _bounded(value: str, limit: int) -> tuple[str, bool]:
    return (value[:limit], True) if len(value) > limit else (value, False)


def _pdf(path: Path, limit: int) -> ExtractedText:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    empty: list[int] = []
    for number, page in enumerate(reader.pages[:MAX_PDF_PAGES], start=1):
        try:
            content = page.extract_text() or ""
        except Exception:
            content = ""
        if content.strip():
            parts.append(f"[página {number}]\n{content.strip()}")
        else:
            empty.append(number)
    text, truncated = _bounded("\n\n".join(parts), limit)
    return ExtractedText(text, truncated, tuple(empty))


def _docx(path: Path, limit: int) -> ExtractedText:
    from docx import Document

    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text, truncated = _bounded("\n".join(parts), limit)
    return ExtractedText(text, truncated)


def _xlsx(path: Path, limit: int) -> ExtractedText:
    from openpyxl import load_workbook

    book = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in book.worksheets:
        parts.append(f"[planilha {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if value is None else str(value) for value in row]
            if any(cell.strip() for cell in cells):
                parts.append(" | ".join(cells))
    book.close()
    text, truncated = _bounded("\n".join(parts), limit)
    return ExtractedText(text, truncated)


def _pptx(path: Path, limit: int) -> ExtractedText:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[slide {number}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                parts.append(text.strip())
    text, truncated = _bounded("\n".join(parts), limit)
    return ExtractedText(text, truncated)


def _plain(path: Path, limit: int) -> ExtractedText:
    data = path.read_bytes()[: limit * 4 + 1]
    text, truncated = _bounded(data.decode("utf-8", "replace"), limit)
    return ExtractedText(text, truncated)


_HANDLERS = {
    "application/pdf": _pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _pptx,
}


def extract_text(path: Path, media_type: str, *, max_chars: int = MAX_EXTRACTED_CHARS) -> ExtractedText:
    """Extract text from a document. Raises ``ValueError`` for image types."""
    handler = _HANDLERS.get(media_type)
    if handler is not None:
        return handler(path, max_chars)
    if media_type.startswith("text/") or media_type in {"application/json", "application/xml"}:
        return _plain(path, max_chars)
    raise ValueError(f"{media_type} has no text to extract")


__all__ = ["ExtractedText", "MAX_EXTRACTED_CHARS", "extract_text"]
